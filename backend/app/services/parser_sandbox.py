"""L-8: parse untrusted uploaded documents in a resource-capped child process.

A malicious file (decompression bomb, pathological allocation, runaway loop in a
parser library) should not be able to exhaust the request worker's memory or CPU.
Each parse runs in a fresh subprocess with an address-space and CPU-second cap.

Limits are POSIX-only (the production container is Linux). On other platforms the
parse still runs out-of-process — isolating a crash — but without rlimits.

Parser libraries are imported lazily inside the worker functions so the spawned
child re-imports only this light module, not the heavy service layer.
"""

from __future__ import annotations

import io
import multiprocessing as mp
import os
from concurrent.futures import ProcessPoolExecutor
from typing import TypeVar

T = TypeVar("T")

# Per-document caps for the sandboxed parse.
_RLIMIT_AS_BYTES = 2 * 1024 * 1024 * 1024  # 2 GiB address space
_RLIMIT_CPU_SECONDS = 120

# Spawn (not fork) so the child does not inherit the parent's threads/locks.
_MP_CONTEXT = mp.get_context("spawn")


def _apply_limits() -> None:
    """Lower the child's resource ceilings. No-op off POSIX."""
    if os.name != "posix":
        return
    import resource

    for res, cap in (
        (resource.RLIMIT_AS, _RLIMIT_AS_BYTES),
        (resource.RLIMIT_CPU, _RLIMIT_CPU_SECONDS),
    ):
        try:
            _soft, hard = resource.getrlimit(res)
            ceiling = cap if hard == resource.RLIM_INFINITY else min(cap, hard)
            resource.setrlimit(res, (ceiling, ceiling))
        except (ValueError, OSError):
            # Some sandboxes forbid raising/lowering a given limit; skip it
            # rather than failing the whole parse.
            continue


def run_sandboxed(func, *args):
    """Run ``func(*args)`` in a fresh, resource-limited worker process.

    ``max_tasks_per_child=1`` guarantees a brand-new process per call so the
    CPU-second cap applies per document rather than accumulating.
    """
    with ProcessPoolExecutor(
        max_workers=1,
        max_tasks_per_child=1,
        mp_context=_MP_CONTEXT,
        initializer=_apply_limits,
    ) as pool:
        return pool.submit(func, *args).result()


# --- Pure, picklable parse functions (module-level so spawn can import them) ---


def extract_text_local(data: bytes, content_type: str) -> str:
    """Extract text from supported formats. No network — OCR fallback for
    text-less PDFs is handled by the caller, in-process."""
    if content_type == "application/pdf":
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data))
        pages: list[str] = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        return "\n\n".join(pages)

    if content_type == (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        from docx import Document as DocxDocument

        doc = DocxDocument(io.BytesIO(data))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())

    if content_type in ("text/plain", "text/csv", "application/msword"):
        return data.decode("utf-8", errors="replace")

    if content_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        rows: list[str] = []
        for sheet in wb.worksheets:
            rows.append(f"--- {sheet.title} ---")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append("\t".join(cells))
        wb.close()
        return "\n".join(rows)

    if content_type == (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        slides: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            parts = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
            slides.append("\n".join(parts))
        return "\n\n".join(slides)

    raise ValueError(f"Unsupported content type: {content_type}")


def render_pdf_pages(data: bytes) -> list[bytes]:
    """Render each PDF page to a PNG via PyMuPDF (parses untrusted bytes)."""
    import fitz  # pymupdf

    doc = fitz.open(stream=data, filetype="pdf")
    images: list[bytes] = []
    for page in doc:
        pix = page.get_pixmap(dpi=150)
        images.append(pix.tobytes("png"))
    doc.close()
    return images
