import asyncio
import base64
import io
import uuid

import structlog
import tiktoken
from azure.ai.documentintelligence import DocumentIntelligenceClient
from azure.ai.documentintelligence.models import AnalyzeDocumentRequest, DocumentContentFormat
from azure.core.credentials import AzureKeyCredential

from sqlalchemy import select

from app.core.config import settings
from app.models.document import DocumentStatus
from app.models.notification import NotificationType
from app.models.user import User
from app.repositories.document import DocumentRepository
from app.services.notification import NotificationDispatcher
from app.services.openai_client import AzureOpenAIClient
from app.services.search_indexer import SearchIndexer
from app.services.storage import BlobStorageService

logger = structlog.get_logger(__name__)


class DocumentProcessor:
    def __init__(
        self,
        storage: BlobStorageService,
        openai_client: AzureOpenAIClient,
        indexer: SearchIndexer,
        repo: DocumentRepository,
    ) -> None:
        self._storage = storage
        self._openai = openai_client
        self._indexer = indexer
        self._repo = repo

    @staticmethod
    def _extract_text_pdf(data: bytes) -> str:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data))
        pages: list[str] = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        return "\n\n".join(pages)

    @staticmethod
    def _extract_text_docx(data: bytes) -> str:
        from docx import Document as DocxDocument

        doc = DocxDocument(io.BytesIO(data))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())

    @staticmethod
    def _extract_text_plain(data: bytes) -> str:
        return data.decode("utf-8", errors="replace")

    @staticmethod
    def _extract_text_xlsx(data: bytes) -> str:
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        rows: list[str] = []
        for sheet in wb.worksheets:
            rows.append(f"--- {sheet.title} ---")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append("\t".join(cells))
        wb.close()
        return "\n".join(rows)

    @staticmethod
    def _extract_text_pptx(data: bytes) -> str:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        slides: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            parts = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
            slides.append("\n".join(parts))
        return "\n\n".join(slides)

    @staticmethod
    def _extract_text_ocr(data: bytes) -> str:
        """Use Azure AI Document Intelligence to OCR a scanned PDF."""
        if not settings.azure_doc_intelligence_endpoint:
            raise ValueError(
                "No text extracted and Azure Document Intelligence is not configured for OCR"
            )

        client = DocumentIntelligenceClient(
            endpoint=settings.azure_doc_intelligence_endpoint,
            credential=AzureKeyCredential(settings.azure_doc_intelligence_key),
        )

        poller = client.begin_analyze_document(
            "prebuilt-read",
            AnalyzeDocumentRequest(bytes_source=data),
            output_content_format=DocumentContentFormat.MARKDOWN,
        )
        result = poller.result()
        return result.content or ""

    @staticmethod
    def _extract_text(data: bytes, content_type: str) -> str:
        if content_type == "application/pdf":
            text = DocumentProcessor._extract_text_pdf(data)
            if not text.strip():
                logger.info("pdf_no_text_layer_falling_back_to_ocr")
                text = DocumentProcessor._extract_text_ocr(data)
            return text
        elif content_type == (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ):
            return DocumentProcessor._extract_text_docx(data)
        elif content_type in ("text/plain", "text/csv"):
            return DocumentProcessor._extract_text_plain(data)
        elif content_type == "application/msword":
            # Old .doc format — attempt plain text extraction as fallback
            return DocumentProcessor._extract_text_plain(data)
        elif content_type == ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"):
            return DocumentProcessor._extract_text_xlsx(data)
        elif content_type == (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ):
            return DocumentProcessor._extract_text_pptx(data)
        else:
            raise ValueError(f"Unsupported content type: {content_type}")

    @staticmethod
    def _render_pdf_pages(data: bytes) -> list[bytes]:
        """Render each PDF page as a PNG image using PyMuPDF."""
        import fitz  # pymupdf

        doc = fitz.open(stream=data, filetype="pdf")
        images: list[bytes] = []
        for page in doc:
            # Render at 150 DPI for good quality without huge size
            pix = page.get_pixmap(dpi=150)
            images.append(pix.tobytes("png"))
        doc.close()
        return images

    async def _describe_pdf_pages(self, data: bytes) -> str:
        """Render PDF pages to images and describe visual content via GPT-4o vision."""
        page_images = await asyncio.to_thread(self._render_pdf_pages, data)
        if not page_images:
            return ""

        prompt = (
            "Describe all visual elements on this page: images, photos, diagrams, charts, "
            "maps, logos, tables, and any other non-text content. Be specific about what "
            "you see (e.g., aircraft type, airport layout features, risk matrix values). "
            "If there are no visual elements beyond plain text, respond with NONE."
        )

        descriptions: list[str] = []
        for i, img_bytes in enumerate(page_images, 1):
            try:
                img_b64 = base64.b64encode(img_bytes).decode("ascii")
                description = await self._openai.describe_image(img_b64, prompt)
                if description.strip().upper() != "NONE":
                    descriptions.append(f"[Page {i} visual content]: {description}")
            except Exception:
                logger.warning("page_vision_failed", page=i, exc_info=True)
                continue

        return "\n\n".join(descriptions)

    @staticmethod
    def _chunk_text(text: str) -> list[str]:
        encoding = tiktoken.get_encoding("cl100k_base")
        tokens = encoding.encode(text)
        chunks: list[str] = []

        start = 0
        while start < len(tokens):
            end = start + settings.chunk_size_tokens
            chunk_tokens = tokens[start:end]
            chunk_text = encoding.decode(chunk_tokens)
            chunks.append(chunk_text)
            start = end - settings.chunk_overlap_tokens

        return chunks

    @staticmethod
    async def _extract_text_in_thread(data: bytes, content_type: str) -> str:
        return await asyncio.to_thread(DocumentProcessor._extract_text, data, content_type)

    @staticmethod
    async def _chunk_text_in_thread(text: str) -> list[str]:
        return await asyncio.to_thread(DocumentProcessor._chunk_text, text)

    async def process(self, document_id: uuid.UUID) -> None:
        await self._repo.update_status(document_id, DocumentStatus.PROCESSING)

        try:
            document = await self._repo.get_by_id_system(document_id)
            if not document:
                logger.error("document_not_found", document_id=str(document_id))
                return

            data = await self._storage.download(document.blob_path)
            text = await DocumentProcessor._extract_text_in_thread(data, document.content_type)

            # For PDFs, analyze pages with GPT-4o vision to capture images/diagrams
            if document.content_type == "application/pdf":
                try:
                    image_descriptions = await self._describe_pdf_pages(data)
                    if image_descriptions:
                        text = text + "\n\n" + image_descriptions
                        logger.info(
                            "vision_descriptions_added",
                            document_id=str(document_id),
                        )
                except Exception:
                    logger.warning(
                        "vision_analysis_failed_continuing",
                        document_id=str(document_id),
                        exc_info=True,
                    )

            if not text.strip():
                await self._repo.update_status(
                    document_id, DocumentStatus.FAILED, error_message="No text extracted"
                )
                return

            chunks = await DocumentProcessor._chunk_text_in_thread(text)
            logger.info(
                "document_chunked",
                document_id=str(document_id),
                chunk_count=len(chunks),
            )

            embeddings: list[list[float]] = []
            batch_size = settings.embedding_batch_size
            for i in range(0, len(chunks), batch_size):
                batch = chunks[i : i + batch_size]
                batch_embeddings = await self._openai.embed_batch(batch)
                embeddings.extend(batch_embeddings)

            await self._indexer.index_chunks(
                document_id=document_id,
                organization_id=document.organization_id,
                source=document.filename,
                chunks=chunks,
                embeddings=embeddings,
                source_type=document.source_type.value,
            )

            await self._repo.update_status(
                document_id, DocumentStatus.INDEXED, chunk_count=len(chunks)
            )
            logger.info(
                "document_processed",
                document_id=str(document_id),
                chunks_indexed=len(chunks),
            )

            uploader_row = await self._repo._db.execute(  # type: ignore[attr-defined]
                select(User).where(User.id == document.uploaded_by)
            )
            uploader = uploader_row.scalar_one_or_none()
            if uploader is not None:
                NotificationDispatcher().dispatch(
                    organization_id=document.organization_id,
                    triggered_by=uploader,
                    notification_type=NotificationType.DOCUMENT_INDEXED,
                    title=f"Document indexed: {document.filename[:100]}",
                    body=f"{len(chunks)} chunks indexed. Source: {document.source_type.value}",
                    resource_type="document",
                    resource_id=str(document.id),
                )

        except Exception as e:
            logger.error(
                "document_processing_failed",
                document_id=str(document_id),
                error=str(e),
                exc_info=True,
            )
            await self._repo.update_status(
                document_id,
                DocumentStatus.FAILED,
                error_message="Document processing failed. See server logs for details.",
            )
